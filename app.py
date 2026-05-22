import re
from io import BytesIO
from collections import Counter
from datetime import date, datetime, timedelta
import calendar

import pandas as pd
import streamlit as st
from dateutil.parser import parse as parse_date
from pptx import Presentation
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter


# --- extraction engine ---
# =========================================================================== #
#  LAYOUT DETECTION
#  --------------------------------------------------------------------------
#  Different law firms / clients use different docket conventions, e.g.
#      Antiva        01394-0005-00EP        signature  #-#-#@
#      Eradivir      2018-LOW-68327-04      signature  #-@-#-#
#      NewAmsterdam  P6046729US1            signature  @#@#
#  but each deck is internally consistent. Rather than hard-coding every
#  client's format, we LEARN this deck's docket pattern in a first pass:
#  abstract each box's leading token into a structural signature (digit-run ->
#  '#', letter-run -> '@', separators kept), tally them, and the dominant
#  docket-shaped signature(s) ARE this deck's docket pattern.
#
#  Only the docket varies between clients. Application / PCT / WIPO numbers are
#  keyed to jurisdiction (a US serial is ##/###,### everywhere), so those
#  patterns -- defined later in this file -- need no per-deck detection.
# =========================================================================== #

_DATE_TOKEN_RE = re.compile(r"^\d{1,2}[/-]\d{1,2}[/-]\d{2,4}$")
_STATUS_WORDS = {"ISSUED", "GRANTED", "ALLOWED", "ABANDONED", "WITHDRAWN",
                 "PENDING", "PRIORITY", "ACTIVE", "CASES", "STRUCTURE",
                 "STRUCTURES", "APPENDIX", "KEY", "PRIORITIES"}


def _signature(token):
    """Abstract a token: digit-run -> '#', letter-run -> '@', separators kept."""
    return re.sub(r"[A-Za-z]+", "@", re.sub(r"\d+", "#", token))


def _first_token(box_text):
    lines = box_text.splitlines()
    if not lines:
        return ""
    parts = lines[0].strip().split()
    return parts[0] if parts else ""


def _is_docket_ish(tok):
    """Could this leading token plausibly be a docket (not a date/app#/word)?"""
    if not tok or len(tok) < 6:
        return False
    if _DATE_TOKEN_RE.match(tok):
        return False
    u = tok.upper()
    if u.startswith(("PCT/", "WO")):
        return False
    if u.strip(",/.") in _STATUS_WORDS:
        return False
    if not re.search(r"\d", tok):                 # a docket always has a digit
        return False
    # reject a bare US-style application serial, e.g. 17/263,451
    if re.fullmatch(r"\d{1,3}/\d{2,3},?\d{0,3}", tok):
        return False
    return True


def _sig_is_docket_shaped(sig):
    """A docket signature mixes letters+digits, or is a multi-group number."""
    has_at = "@" in sig
    has_hash = "#" in sig
    groups = len(re.split(r"[-/._]", sig))
    if has_at and has_hash and len(sig) >= 3:
        return True
    if has_hash and groups >= 3 and len(sig) >= 5:
        return True
    return False


def _sig_to_regex(sig):
    """Concrete regex for a signature: '#' -> \\d+, '@' -> [A-Za-z]+, sep literal."""
    out = []
    for ch in sig:
        if ch == "#":
            out.append(r"\d+")
        elif ch == "@":
            out.append(r"[A-Za-z]+")
        else:
            out.append(re.escape(ch))
    return "".join(out)


def detect_docket_regex(box_texts):
    """
    Learn this deck's docket pattern from its leading tokens.

    Returns (compiled_regex, signature_list). The regex matches a docket at the
    START of a token. Returns (None, []) when no docket-shaped signature recurs,
    in which case the caller falls back to the built-in patterns.
    """
    cands = [_first_token(b) for b in box_texts]
    cands = [c for c in cands if _is_docket_ish(c)]
    if not cands:
        return None, []

    sigs = Counter(_signature(c) for c in cands)
    threshold = max(3, int(0.04 * len(box_texts)))
    keep = [s for s, n in sigs.items()
            if n >= threshold and _sig_is_docket_shaped(s)]
    if not keep:
        # nothing cleared the bar; take the most common docket-shaped signature
        keep = [s for s, _ in sigs.most_common() if _sig_is_docket_shaped(s)][:1]
    if not keep:
        return None, []

    # Longest signatures first so a fuller docket wins over a prefix of itself
    # (e.g. P6046729PCT-CN before P6046729EP's @#@).
    keep.sort(key=len, reverse=True)
    pattern = "|".join(_sig_to_regex(s) for s in keep)
    return re.compile(r"^(?:%s)" % pattern), keep


def detect_matter_regex(docket_signatures):
    """
    Derive a 'matter family' prefix regex from the learned docket signatures,
    used to group related cases (all national filings off one matter).

    The matter is the docket minus its trailing case suffix. We strip a trailing
    country+sequence (e.g. 'US1', '-00EP', '-04', 'PCT-CN') heuristically.
    Returns a compiled regex whose group(1) is the matter stem, or None.
    """
    if not docket_signatures:
        return None
    # Build from the longest signature: keep everything up to the last
    # letter/sep transition that looks like a case suffix.
    # Practical approach: matter = docket with a trailing
    #   (separator? + letters + digits?)  OR  (separator + digits)  removed.
    return re.compile(
        r"^(.*?)(?:[-/]?[A-Za-z]{2,}[-/]?[A-Za-z0-9]*\d*|[-/]\d{1,3})$"
    )


# --------------------------------------------------------------------------- #
#  Docket / client-code patterns
# --------------------------------------------------------------------------- #

# Two docket conventions are supported:
#   Antiva-style:   01394-0005-00EP   /  01394-0003-00MO-01CN
#                   (5 digits - 4 digits - 2 digits + country, optional 2nd leg)
#   Eradivir-style: 2018-LOW-68327-04
#                   (4-digit year - 3-letter client mnemonic - 4-6 digit matter
#                    - 2-digit case suffix; country is NOT in the docket here,
#                    it follows the application number on the same/next line)
DOCKET_RE = re.compile(
    r"\b\d{5}-\d{4}-\d{2}[A-Z]{2,4}(?:-\d{2}[A-Z]{2,4})?\b"     # Antiva
    r"|\b\d{4}-[A-Z]{2,4}-\d{4,6}-\d{2}\b"                      # Eradivir
)

# Client code tokens that we want to ignore when hunting for an app number:
#   15080-101EP1, 15080-004US3-CON2, 15080-101WO1US2, 15040-004CN2-MO, 107USP1
#
# Key constraint: a client code is <4-5 digits>-<3 digits> followed by a LETTER
# (the country/type segment). Requiring that trailing letter is what stops the
# pattern from swallowing pure-numeric application numbers such as the Japanese
# "2016-502307" or the Korean "...-2025-7005473".
CLIENT_CODE_RES = [
    re.compile(r"\b\d{4,5}-\d{3}[A-Za-z][A-Za-z0-9-]*"),   # 15080-101EP1, 15040-004CN2-MO
    re.compile(r"\b\d{3}[A-Z]{2,3}P?\d*\b"),               # 107USP1, 107USP2
]


# --------------------------------------------------------------------------- #
#  Country-specific APPLICATION-number patterns (validated against the deck)
# --------------------------------------------------------------------------- #
COUNTRY_APP_RES = {
    "US": re.compile(r"\b\d{2}/\d{3},?\d{3}\b"),          # 61/793,993 ; 14/211,002
    "EP": re.compile(r"\b\d{8}\.\d\b"),                   # 14764430.6
    "JP": re.compile(r"\b\d{4}-\d{6}\b"),                 # 2016-502307
    "KR": re.compile(r"\b10-\d{4}-\d{7}\b"),              # 10-2017-7008850
    "CN": re.compile(r"(?<!\d)\d{12}\.\s?[\dX](?!\d)"),   # 201580054350.9 / .X / ZL-prefixed
    "IN": re.compile(r"\b\d{12}\b"),                      # 201747008733
    "AU": re.compile(r"\b\d{10}\b"),                      # 2015317972
    "TW": re.compile(r"\b\d{9}\b"),                       # 112127315
    "CL": re.compile(r"\b\d{9}\b"),                       # 202500166
    "EA": re.compile(r"\b\d{9}\b"),                       # 201790628
    "CA": re.compile(r"\b\d,?\d{3},?\d{3}\b"),            # 2,961,200 / 3262284
    "IL": re.compile(r"\b\d{6}\b"),                       # 250836
    "NZ": re.compile(r"\b\d{6,7}\b"),                     # 818923
    "SG": re.compile(r"\b\d{11}[A-Z]\b"),                 # 11201701957X
    "ZA": re.compile(r"\b\d{4}/\d{5}\b"),                 # 2025/01545
    "MX": re.compile(r"\bMX/a/\d{4}/\d{6}\b"),            # MX/a/2017/002789
    "CO": re.compile(r"\bNC\d{4}/\d{6}\b"),               # NC2025/001854
    "AR": re.compile(r"\bP\d{9}\b"),                      # P230101911
    "HK": re.compile(r"(?<!\d)\d{8,11}\.\d(?!\d)"),       # 17113734.5 / 62024096696.5
    "BR": re.compile(r"BR[\d\s]+?\.\d"),                  # BR112017005111.7 / BR 12 2022 023284.1
    "MO": re.compile(r"\bJ/\d+\b"),                       # J/008517 (Macau)
}

# Order used for the "try everything" fallback when the docket country is
# unknown or its pattern misses. More-specific patterns come first so they win.
FALLBACK_ORDER = [
    "MX", "CO", "AR", "BR", "KR", "ZA", "SG", "US", "MO",
    "CN", "HK", "EP", "JP", "IN", "AU", "TW", "CA", "IL", "NZ",
]

PCT_RE = re.compile(r"\bPCT/[A-Z]{2}\d{4}/\d{5,6}\b")
# WO2014/143643 ; WO2024020127 (no slash) ; WO 2025/160227 A1
WIPO_RE = re.compile(r"\bWO\s?\d{4}/?\d{6}(?:\s?A\d)?\b")

# US grant numbers like 9,629,860 / 10,195,222 (used for the bonus "patent no." column)
US_GRANT_RE = re.compile(r"\b\d{1,2},\d{3},\d{3}\b")

DATE_RE = re.compile(r"\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b")
DUE_LINE_RE = re.compile(r"\b(due|by)\b", re.IGNORECASE)

STATUS_RE = re.compile(r"\b(ABN|ABANDONED|WITHDRAWN|PENDING|GRANTED|ISSUED|EXPIRED)\b",
                       re.IGNORECASE)

# Boxes that are pure headers / page chrome (no real case data)
SKIP_EXACT = {"structure", "case structures", "key", "priorities",
              "appendix \u2013 blow-ups", "appendix - blow-ups"}


# --------------------------------------------------------------------------- #
#  Helpers
# --------------------------------------------------------------------------- #
def _country_from_docket_suffix(docket):
    """Antiva style: EP from 01394-0005-00EP ; MO special-cased."""
    if not docket:
        return None
    if "MO" in docket:
        return "MO"
    last = docket.split("-")[-1]
    m = re.match(r"\d*([A-Z]{2,4})$", last)
    return m.group(1) if m else None


# Known jurisdiction codes that appear as a trailing 2-letter code after the
# application number (Eradivir style: "17/263,451  US", "3107778  CA").
KNOWN_COUNTRIES = {
    "US", "EP", "EU", "JP", "KR", "CN", "IN", "AU", "TW", "CL", "EA", "CA",
    "IL", "NZ", "SG", "ZA", "MX", "CO", "AR", "HK", "BR", "MO", "MY", "PH",
    "TH", "VN", "ID", "GB", "DE", "FR", "PE", "EC", "CR", "GT",
    # Additional jurisdictions seen across client decks:
    "RU", "UA", "DZ", "EG", "LY", "TN", "PA", "AP", "VE", "NL", "AT", "SA",
    "AE", "KW", "UY", "LB", "OA", "GC", "QA", "BH", "JO", "MA", "NG", "PK",
    "BD", "LK", "KZ", "AZ", "GE", "RS", "UA",
}


def _country_from_app_line(text):
    """
    The 2-letter jurisdiction follows the application number, e.g. a line
    ending in '...  US' or '... CN'. It may carry a trailing continuation /
    status tag that we ignore:
        '19/257,226  USC4'   -> US   (continuation 4)
        '63/891,043  US P2'  -> US   (provisional 2)
        '62025115409.7 HK-CN'-> HK   (HK based on a CN parent)
        '...  EP D1'         -> EP   (divisional 1)
    Returns the first jurisdiction code found.
    """
    for line in text.splitlines():
        line = line.rstrip()
        # <app-number char> <space> <CC> [optional tag: space/hyphen + letters/digits]
        m = re.search(
            r"[\dA-Za-z,./)\-]\s+([A-Z]{2})(?:[\s\-]?[A-Z]{0,2}\d{0,2}|\d{1,2})?\s*$",
            line,
        )
        if m and m.group(1) in KNOWN_COUNTRIES:
            return m.group(1)
    return None


def get_country(docket, text=""):
    """
    Resolve the jurisdiction.

    Priority:
      1) Trailing 2-letter code after the application number (Eradivir/NewAms).
      2) Trailing letters of the docket suffix (Antiva decks).
      3) A bare US-style serial (##/###,### or 6x/...) with no country code:
         infer US, since that format is unambiguously a US application.
    Normalizes EU -> EP since both denote a European application here.
    """
    country = _country_from_app_line(text) or _country_from_docket_suffix(docket)
    if country is None and re.search(r"\b\d{2}/\d{3},\d{3}\b", text):
        country = "US"
    if country == "EU":
        country = "EP"
    return country


def _strip_known_tokens(text, docket):
    """Remove docket + client-code tokens so they don't masquerade as app #s."""
    out = text
    if docket:
        out = out.replace(docket, " ")
    for rx in CLIENT_CODE_RES:
        out = rx.sub(" ", out)
    return out


def _clean_match(value):
    if value is None:
        return None
    # collapse internal whitespace (e.g. CN "201580054350. 9" -> "201580054350.9")
    return re.sub(r"\s+", "", value.strip())


def find_application_number(text, docket):
    """Return (application_number, country)."""
    search_text = _strip_known_tokens(text, docket)
    country = get_country(docket, text)

    # PCT national-phase parent: the PCT number is the identifier, not a
    # separate application serial. Reported in the PCT column instead.
    if country == "PCT":
        return None, country

    # A box whose only identifier is a PCT number (no national jurisdiction
    # code present) is a PCT filing -- don't coin a serial from the PCT digits.
    if country is None and PCT_RE.search(text):
        return None, "PCT"

    # 1) Use the jurisdiction's own pattern. For a known country we trust only
    #    that pattern -- guessing with another country's (looser) pattern tends
    #    to grab grant numbers, so we'd rather leave the field blank.
    if country in COUNTRY_APP_RES:
        m = COUNTRY_APP_RES[country].search(search_text)
        if m:
            return _clean_match(m.group(0)), country
        # Known country but its pattern missed (e.g. a typo'd serial in the
        # source slide). Fall through to the last-resort capture below rather
        # than dropping the number entirely.

    # 2) Unknown jurisdiction -> try every pattern in priority order.
    if country not in COUNTRY_APP_RES:
        for key in FALLBACK_ORDER:
            m = COUNTRY_APP_RES[key].search(search_text)
            if m:
                return _clean_match(m.group(0)), country

    # 3) Country-anchored generic capture: take the token sitting right before a
    #    trailing 2-letter country code. Handles jurisdictions we have no
    #    explicit pattern for (DZ, AP, TH, VN, UY, ...).
    anchored = _appno_before_country(text, docket)
    if anchored:
        return anchored, country

    # 4) Last resort: capture the most number-like token on the line directly
    #    after the docket, so a malformed entry is surfaced (flagged for review)
    #    instead of silently disappearing.
    raw = _last_resort_appno(text, docket)
    return raw, country


def _appno_before_country(text, docket):
    """
    Generic capture: the application number is the token immediately before a
    trailing 2-letter country code, on any line. Works for jurisdictions we
    don't have an explicit pattern for ('P6046729PCT-DZ 10734 DZ' -> 10734,
    'I651086B  TW' -> I651086B).
    """
    for line in text.splitlines():
        line = line.rstrip()
        m = re.search(r"([A-Za-z]{0,2}\d[\dA-Za-z,./\-]*)\s+([A-Z]{2})\s*$", line)
        if m and m.group(2) in KNOWN_COUNTRIES:
            token = m.group(1)
            if docket and token in docket:
                continue
            if re.fullmatch(r"(?:19|20)\d{2}", token):   # a real year, not an app #
                continue
            return _clean_match(token)
    return None


def _last_resort_appno(text, docket):
    """Grab a plausible application-number token when no pattern matched."""
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    # Find the docket line, then look at it and the following line.
    candidates = []
    for i, ln in enumerate(lines):
        if docket and docket in ln:
            tail = ln.replace(docket, " ")
            candidates.append(tail)
            if i + 1 < len(lines):
                candidates.append(lines[i + 1])
            break
    for cand in candidates:
        # a token with at least 5 digits and a slash/comma/dot, e.g. 674/460,651
        m = re.search(r"\b[\dA-Z]{0,3}[\d][\d,./\-]{4,}\d\b", cand)
        if m:
            return _clean_match(m.group(0))
    return None


def find_pct(text):
    m = PCT_RE.search(text)
    return m.group(0) if m else None


def find_wipo(text):
    m = WIPO_RE.search(text)
    if not m:
        return None
    # Normalize to "WO YYYY/NNNNNN" (keep optional kind code like A1)
    raw = re.sub(r"\s+", "", m.group(0))           # WO2025/160227A1
    raw = re.sub(r"^WO", "", raw)
    kind = ""
    km = re.search(r"(A\d)$", raw)
    if km:
        kind = " " + km.group(1)
        raw = raw[: km.start()]
    return f"WO {raw}{kind}".strip()


def _norm_date(raw):
    """Return MM/DD/YYYY or None."""
    try:
        d = parse_date(raw, dayfirst=False, fuzzy=False)
        return d.strftime("%m/%d/%Y")
    except Exception:
        return None


def find_dates(lines):
    """
    Split dates into a single filing date and a list of due-date deadlines.

    * A *due date* is any date that sits on a line containing 'due'/'by'.
    * The *filing date* is the first non-due date (typically next to the app #).
    Returns (filing_date, [ {action, date}, ... ], [undated_action, ...]).
    """
    filing = None
    deadlines = []
    undated = []

    for line in lines:
        line = line.strip()
        if not line:
            continue
        dates_on_line = DATE_RE.findall(line)
        is_due = bool(DUE_LINE_RE.search(line))

        if is_due:
            if dates_on_line:
                for raw in dates_on_line:
                    nd = _norm_date(raw)
                    if not nd:
                        continue
                    # action = text up to the date, tidied
                    idx = line.find(raw)
                    action = line[:idx].strip(" :-\u2013") or line.strip()
                    deadlines.append({"action": action, "date": nd})
            else:
                # Undated action item, e.g. "Assignment due".
                # Require the explicit word "due" so a stray "...by..." with no
                # date (common in narrative notes) isn't logged as an action.
                if re.search(r"\bdue\b", line, re.IGNORECASE):
                    undated.append(line)
        else:
            for raw in dates_on_line:
                nd = _norm_date(raw)
                if nd and filing is None:
                    filing = nd
                    break

    return filing, deadlines, undated


# --------------------------------------------------------------------------- #
#  Box-level parsing
# --------------------------------------------------------------------------- #
def parse_box(text, slide_num, docket_re=None):
    """
    Return a case dict for one text box, or None if it holds no case data.

    docket_re: the per-deck docket pattern learned by detect_docket_regex().
    When it doesn't match (or isn't supplied), fall back to the built-in
    DOCKET_RE so previously-supported decks keep working.
    """
    raw = text.strip()
    if raw.lower() in SKIP_EXACT:
        return None

    lines = [ln.strip() for ln in raw.splitlines() if ln.strip()]
    if not lines:
        return None

    docket = None
    # 1) Learned pattern, anchored to the box's first token.
    if docket_re is not None:
        ft = _first_token(raw)
        m = docket_re.match(ft)
        if m:
            docket = m.group(0)
    # 2) Fallback: the built-in multi-format pattern, searched anywhere.
    if docket is None:
        m = DOCKET_RE.search(raw)
        if m:
            docket = m.group(0)

    pct = find_pct(raw)
    wipo = find_wipo(raw)

    # Skip boxes that have no identifiers at all (titles, page numbers, etc.)
    if not docket and not pct and not wipo:
        return None

    app_no, country = (None, None)
    if docket:
        app_no, country = find_application_number(raw, docket)
        # Single-identifier box like "I338684B  TW": the detector treated the
        # number as a docket, but it's really the application/grant number with
        # no separate docket. Re-assign so the number isn't lost.
        if app_no is None and country and len(lines) <= 2:
            m = re.match(r"^(\S+)\s+([A-Z]{2})\s*$", lines[0])
            if m and m.group(1) == docket and m.group(2) in KNOWN_COUNTRIES:
                app_no, docket = docket, None
    elif pct:
        country = "PCT"

    filing, deadlines, undated = find_dates(lines)

    status_m = STATUS_RE.search(raw)
    status = status_m.group(0).upper() if status_m else ""

    return {
        "slide": slide_num,
        "docket": docket,
        "country": country,
        "application_number": app_no,
        "pct_number": pct,
        "wipo_number": wipo,
        "filing_date": filing,
        "status": status,
        "deadlines": deadlines,
        "undated_actions": undated,
        "raw_text": raw,
    }


# --------------------------------------------------------------------------- #
#  Shape walking
# --------------------------------------------------------------------------- #
def iter_box_texts(shape):
    """Yield text for a shape, recursing into groups."""
    if shape.shape_type == 6:  # GroupShape
        for child in shape.shapes:
            yield from iter_box_texts(child)
    elif getattr(shape, "has_text_frame", False) and shape.text.strip():
        yield shape.text.strip()


def extract_cases(pptx_source):
    """Parse a pptx (path or file-like) -> list of case dicts."""
    prs = Presentation(pptx_source)

    # First pass: collect every text box so we can learn this deck's docket
    # format before extracting anything.
    box_index = []   # (slide_num, text)
    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            for text in iter_box_texts(shape):
                box_index.append((slide_num, text))

    docket_re, _sigs = detect_docket_regex([t for _, t in box_index])

    # Second pass: parse each box using the learned pattern.
    cases = []
    for slide_num, text in box_index:
        case = parse_box(text, slide_num, docket_re=docket_re)
        if case:
            cases.append(case)
    return cases


# --- excel + calendar builder ---
HEADER_FILL = PatternFill("solid", fgColor="1F3864")
HEADER_FONT = Font(name="Arial", bold=True, color="FFFFFF", size=11)
BODY_FONT = Font(name="Arial", size=10)
THIN = Side(style="thin", color="D9D9D9")
BORDER = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)


def _join_deadlines(case):
    parts = [f"{d['action']}: {d['date']}" for d in case["deadlines"]]
    parts += [u for u in case["undated_actions"]]
    return "; ".join(parts)


def cases_to_rows(cases, client, deadline_cutoff=None):
    """
    Build (deadline_rows, case_rows) as lists of dicts.

    deadline_cutoff: optional datetime.date. When given, deadline rows whose
    due date falls *before* this date are dropped from the Deadlines list.
    The All Cases list is always complete (it's a full inventory, not a
    calendar), so the cutoff never removes a case from it.
    """
    case_rows, deadline_rows, seen = [], [], set()

    for c in cases:
        # Flag rows a human should eyeball: an identifier was captured but the
        # jurisdiction couldn't be resolved (often a typo in the source slide),
        # or a docket has no application/PCT/WIPO number at all.
        review = ""
        if c["docket"] and c["application_number"] and not c["country"]:
            review = "Check: country not resolved"
        elif c["docket"] and not c["application_number"] \
                and not c["pct_number"] and not c["wipo_number"]:
            review = "Check: no application number"

        case_rows.append({
            "Review": review,
            "Client": client,
            "Slide": c["slide"],
            "Docket Number": c["docket"],
            "Country": c["country"],
            "Application Number": c["application_number"],
            "PCT Number": c["pct_number"],
            "WIPO Number": c["wipo_number"],
            "Filing Date": c["filing_date"],
            "Status": c["status"],
            "Due Dates / Actions": _join_deadlines(c),
        })
        for d in c["deadlines"]:
            key = (client, c["docket"], d["action"], d["date"])
            if key in seen:
                continue
            seen.add(key)
            deadline_rows.append({
                "Due Date": d["date"],
                "Action": d["action"],
                "Docket Number": c["docket"],
                "Country": c["country"],
                "Application Number": c["application_number"]
                                      or c["pct_number"] or c["wipo_number"],
                "Client": client,
                "Slide": c["slide"],
            })

    deadline_rows.sort(key=lambda r: datetime.strptime(r["Due Date"], "%m/%d/%Y"))
    if deadline_cutoff is not None:
        deadline_rows = [
            r for r in deadline_rows
            if datetime.strptime(r["Due Date"], "%m/%d/%Y").date() >= deadline_cutoff
        ]
    return deadline_rows, case_rows


def _write_sheet(ws, rows, columns, date_cols=()):
    ws.append(columns)
    for col_idx, _ in enumerate(columns, start=1):
        cell = ws.cell(row=1, column=col_idx)
        cell.fill = HEADER_FILL
        cell.font = HEADER_FONT
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border = BORDER

    for row in rows:
        ws.append([row.get(c, "") for c in columns])

    # body styling + column widths
    widths = {c: len(c) for c in columns}
    for r in range(2, ws.max_row + 1):
        for col_idx, col_name in enumerate(columns, start=1):
            cell = ws.cell(row=r, column=col_idx)
            cell.font = BODY_FONT
            cell.border = BORDER
            cell.alignment = Alignment(vertical="top",
                                       wrap_text=(col_name == "Due Dates / Actions"))
            val = cell.value
            if val not in (None, ""):
                widths[col_name] = max(widths[col_name], min(len(str(val)), 60))

    for col_idx, col_name in enumerate(columns, start=1):
        ws.column_dimensions[get_column_letter(col_idx)].width = widths[col_name] + 2

    ws.freeze_panes = "A2"
    ws.auto_filter.ref = f"A1:{get_column_letter(len(columns))}{ws.max_row}"


def build_workbook(cases, client, deadline_cutoff=None):
    deadline_rows, case_rows = cases_to_rows(cases, client, deadline_cutoff)

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Deadlines"
    _write_sheet(ws1, deadline_rows,
                 ["Due Date", "Action", "Docket Number", "Country",
                  "Application Number", "Client", "Slide"])

    ws2 = wb.create_sheet("All Cases")
    _write_sheet(ws2, case_rows,
                 ["Review", "Client", "Slide", "Docket Number", "Country",
                  "Application Number", "PCT Number", "WIPO Number",
                  "Filing Date", "Status", "Due Dates / Actions"])

    return wb


def workbook_bytes(cases, client, deadline_cutoff=None):
    buf = BytesIO()
    build_workbook(cases, client, deadline_cutoff).save(buf)
    buf.seek(0)
    return buf


# ===========================================================================
#  Calendar outputs:  .ics file  +  in-app month grid
# ===========================================================================

# ===========================================================================
#  Calendar outputs:  printable PDF month  +  in-app HTML month grid
# ===========================================================================

# Brand palette (shared by the PDF and the on-screen grid)
NAVY = (0x1F / 255, 0x38 / 255, 0x64 / 255)
LIGHT = (0xE8 / 255, 0xEE / 255, 0xF7 / 255)
GREY = (0.6, 0.6, 0.6)
RULE = (0x85 / 255, 0x85 / 255, 0x85 / 255)


def month_pdf(deadline_rows, year, month, client_label=""):
    """
    Render a single month as a printable, landscape PDF calendar.

    Returns the PDF as bytes. Pure reportlab (no system libraries), so it
    renders identically on Streamlit Cloud and any desktop.
    """
    from reportlab.lib.pagesizes import letter, landscape
    from reportlab.pdfgen import canvas as _canvas
    from reportlab.lib.units import inch

    # Bucket this month's deadlines by day-of-month.
    events = {}
    for r in deadline_rows:
        d = datetime.strptime(r["Due Date"], "%m/%d/%Y").date()
        if d.year == year and d.month == month:
            events.setdefault(d.day, []).append(r)

    weeks = calendar.Calendar(firstweekday=6).monthdayscalendar(year, month)

    buf = BytesIO()
    page_w, page_h = landscape(letter)            # 792 x 612
    c = _canvas.Canvas(buf, pagesize=(page_w, page_h))

    margin = 0.4 * inch
    title_h = 0.55 * inch
    foot_h = 0.28 * inch
    head_h = 0.24 * inch

    grid_x = margin
    grid_top = page_h - margin - title_h
    grid_w = page_w - 2 * margin
    grid_h = grid_top - margin - foot_h
    col_w = grid_w / 7
    body_top = grid_top - head_h
    row_h = (grid_h - head_h) / len(weeks)

    # --- Title ---
    c.setFillColorRGB(*NAVY)
    c.setFont("Helvetica-Bold", 20)
    c.drawString(margin, page_h - margin - 14, "DocketPoint")
    c.setFont("Helvetica", 14)
    title = f"{calendar.month_name[month]} {year}"
    c.drawRightString(page_w - margin, page_h - margin - 13, title)
    if client_label:
        c.setFillColorRGB(*GREY)
        c.setFont("Helvetica", 9)
        c.drawString(margin, page_h - margin - 30, client_label)

    # --- Weekday header band ---
    c.setFillColorRGB(*NAVY)
    c.rect(grid_x, body_top, grid_w, head_h, fill=1, stroke=0)
    c.setFillColorRGB(1, 1, 1)
    c.setFont("Helvetica-Bold", 9)
    for i, wd in enumerate(_WEEKDAYS):
        c.drawCentredString(grid_x + col_w * (i + 0.5), body_top + 7, wd)

    # --- Day cells ---
    for w, week in enumerate(weeks):
        y_top = body_top - w * row_h
        for i, day in enumerate(week):
            x = grid_x + i * col_w
            # cell border
            c.setStrokeColorRGB(*RULE)
            c.setLineWidth(0.5)
            c.rect(x, y_top - row_h, col_w, row_h, fill=0, stroke=1)
            if day == 0:
                continue
            evs = events.get(day, [])
            # day number
            c.setFont("Helvetica-Bold" if evs else "Helvetica", 9)
            c.setFillColorRGB(*(NAVY if evs else GREY))
            c.drawString(x + 4, y_top - 12, str(day))
            # event pills
            ey = y_top - 24
            for r in evs:
                if ey < y_top - row_h + 6:        # ran out of vertical room
                    c.setFont("Helvetica-Oblique", 6.5)
                    c.setFillColorRGB(*GREY)
                    c.drawString(x + 4, y_top - row_h + 3, "+ more")
                    break
                pill_h = 19
                c.setFillColorRGB(*LIGHT)
                c.rect(x + 3, ey - pill_h + 11, col_w - 6, pill_h, fill=1, stroke=0)
                c.setFillColorRGB(*NAVY)
                c.rect(x + 3, ey - pill_h + 11, 2, pill_h, fill=1, stroke=0)
                c.setFillColorRGB(0.1, 0.1, 0.1)
                c.setFont("Helvetica-Bold", 6.5)
                c.drawString(x + 7, ey + 3, _truncate(r["Docket Number"], col_w - 12, 6.5, bold=True))
                c.setFont("Helvetica", 6.5)
                c.drawString(x + 7, ey - 5,
                             _truncate(r["Action"], col_w - 12, 6.5))
                ey -= pill_h + 3

    # --- Footer ---
    c.setFillColorRGB(*GREY)
    c.setFont("Helvetica", 7)
    stamp = datetime.now().strftime("%m/%d/%Y")
    c.drawString(margin, margin - 2,
                 f"Generated {stamp} \u2022 deadlines extracted from case-structure slides; "
                 f"verify against the system of record.")
    c.drawRightString(page_w - margin, margin - 2,
                      f"{sum(len(v) for v in events.values())} deadline(s) this month")

    c.showPage()
    c.save()
    buf.seek(0)
    return buf.getvalue()


def _truncate(text, max_width_pt, font_size, bold=False):
    """Trim a string with an ellipsis so it fits max_width_pt at font_size."""
    from reportlab.pdfbase.pdfmetrics import stringWidth
    font = "Helvetica-Bold" if bold else "Helvetica"
    text = str(text)
    if stringWidth(text, font, font_size) <= max_width_pt:
        return text
    while text and stringWidth(text + "\u2026", font, font_size) > max_width_pt:
        text = text[:-1]
    return text + "\u2026"


def deadline_months(deadline_rows):
    """Sorted unique (year, month) pairs that contain at least one deadline."""
    months = {(datetime.strptime(r["Due Date"], "%m/%d/%Y").year,
               datetime.strptime(r["Due Date"], "%m/%d/%Y").month)
              for r in deadline_rows}
    return sorted(months)


def month_label(ym):
    """(2026, 7) -> 'July 2026'."""
    return f"{calendar.month_name[ym[1]]} {ym[0]}"


_WEEKDAYS = ["Sun", "Mon", "Tue", "Wed", "Thu", "Fri", "Sat"]


def _html_escape(s):
    return (str(s).replace("&", "&amp;").replace("<", "&lt;")
            .replace(">", "&gt;").replace("'", "&#39;"))


def month_grid_html(deadline_rows, year, month):
    """Return an HTML month calendar with deadlines placed on their due dates."""
    events = {}
    for r in deadline_rows:
        d = datetime.strptime(r["Due Date"], "%m/%d/%Y").date()
        if d.year == year and d.month == month:
            events.setdefault(d.day, []).append(r)

    cal = calendar.Calendar(firstweekday=6)  # Sunday first
    weeks = cal.monthdayscalendar(year, month)

    css = (
        "<style>"
        ".dp-cal{border-collapse:collapse;width:100%;font-family:Arial,sans-serif;"
        "table-layout:fixed}"
        ".dp-cal th{background:#1F3864;color:#fff;padding:6px;font-size:12px;"
        "border:1px solid #d9d9d9}"
        ".dp-cal td{vertical-align:top;height:104px;border:1px solid #d9d9d9;"
        "padding:4px;width:14.28%}"
        ".dp-cal td.empty{background:#fafafa}"
        ".dp-day{font-size:12px;color:#999;text-align:right}"
        ".dp-day.has{color:#1F3864;font-weight:bold}"
        ".dp-ev{background:#E8EEF7;border-left:3px solid #1F3864;border-radius:3px;"
        "padding:2px 5px;margin-top:3px;font-size:11px;line-height:1.25;color:#1a1a1a}"
        ".dp-ev b{display:block;font-size:11px}"
        "</style>"
    )

    html = [css,
            f"<h4 style='font-family:Arial;margin:4px 0 8px'>"
            f"{calendar.month_name[month]} {year}</h4>",
            "<table class='dp-cal'><thead><tr>"]
    html += [f"<th>{w}</th>" for w in _WEEKDAYS]
    html.append("</tr></thead><tbody>")

    for week in weeks:
        html.append("<tr>")
        for day in week:
            if day == 0:
                html.append("<td class='empty'></td>")
                continue
            evs = events.get(day, [])
            day_cls = "dp-day has" if evs else "dp-day"
            cell = [f"<td><div class='{day_cls}'>{day}</div>"]
            for r in evs:
                tip = (f"{r['Action']} \u2014 {r['Docket Number']} "
                       f"({r.get('Country', '')})  App {r.get('Application Number', '')}")
                cell.append(
                    f"<div class='dp-ev' title='{_html_escape(tip)}'>"
                    f"<b>{_html_escape(r['Docket Number'])}</b>"
                    f"{_html_escape(r['Action'])}</div>")
            cell.append("</td>")
            html.append("".join(cell))
        html.append("</tr>")

    html.append("</tbody></table>")
    return "".join(html)


# --- streamlit ui ---
st.set_page_config(page_title="DocketPoint", page_icon="\U0001F4C5",
                   layout="wide", initial_sidebar_state="expanded")

# --- Brand styling -------------------------------------------------------- #
st.markdown("""
<style>
  :root { --dp-navy:#1F3864; --dp-accent:#2F5496; }
  .block-container { padding-top: 2.2rem; max-width: 1200px; }
  #dp-header { display:flex; align-items:center; gap:.7rem; margin-bottom:.15rem; }
  #dp-header .mark {
    width:34px; height:34px; border-radius:8px; background:var(--dp-navy);
    display:flex; align-items:center; justify-content:center;
  }
  #dp-header .mark span { color:#fff; font-size:19px; font-weight:700; }
  #dp-header h1 {
    font-size:30px; font-weight:700; color:var(--dp-navy);
    margin:0; letter-spacing:-.5px;
  }
  #dp-tagline { color:#5a6b85; font-size:14px; margin:0 0 1.1rem 2px; }
  div[data-testid="stFileUploader"] section { border-radius:10px; }
  .stTabs [data-baseweb="tab-list"] { gap: 4px; }
  .stTabs [data-baseweb="tab"] { font-weight:600; }
  div[data-testid="stMetricValue"] { color:var(--dp-navy); }
  .stDownloadButton button {
    background:var(--dp-navy); color:#fff; border:0; border-radius:8px;
    font-weight:600; padding:.5rem 1rem;
  }
  .stDownloadButton button:hover { background:var(--dp-accent); color:#fff; }
</style>
""", unsafe_allow_html=True)

st.markdown("""
<div id="dp-header">
  <div class="mark"><span>DP</span></div>
  <h1>DocketPoint</h1>
</div>
<p id="dp-tagline">Patent docketing, extracted from your case-structure slides.</p>
""", unsafe_allow_html=True)

# --- Sidebar -------------------------------------------------------------- #
with st.sidebar:
    st.markdown("### About")
    st.markdown(
        "DocketPoint reads patent **case-structure** PowerPoint decks and pulls "
        "out each case's docket, application, PCT and WIPO numbers, filing dates, "
        "and dated deadlines."
    )
    st.markdown("### How to use")
    st.markdown(
        "1. Upload one or more `.pptx` decks.\n"
        "2. Review the **Deadlines**, **All Cases**, and **Calendar** tabs.\n"
        "3. Download the Excel workbook or a printable PDF calendar."
    )
    st.markdown("### Lookback")
    st.markdown(
        "By default the calendar and deadline list show only upcoming dates. "
        "Use the slider to also include deadlines that fell due in recent months."
    )
    st.divider()
    st.caption(
        "Deadlines are extracted from slide text and are a convenience view only. "
        "Always verify against the official docketing system of record."
    )

# --- Inputs --------------------------------------------------------------- #
ppt_files = st.file_uploader(
    "Upload case-structure PowerPoint files (.pptx)",
    type="pptx",
    accept_multiple_files=True,
)

months_back = st.slider(
    "Include deadlines due up to this many months in the past",
    0, 36, 0,
    help="0 shows only deadlines from today onward.",
)

if not ppt_files:
    st.info("Upload a case-structure deck (.pptx) to begin.")
    st.stop()

cutoff = date.today() - timedelta(days=int(30.4 * months_back))

all_cases = []
for f in ppt_files:
    client = f.name.replace(".pptx", "")
    cases = extract_cases(f)
    if not cases:
        st.warning(f"No extractable cases found in {f.name}.")
        continue
    all_cases.append((client, cases))

if not all_cases:
    st.stop()

# Apply the cutoff at the source so the tables, calendar, and downloads all agree.
deadline_rows, case_rows = [], []
for client, cases in all_cases:
    d_rows, c_rows = cases_to_rows(cases, client, deadline_cutoff=cutoff)
    deadline_rows += d_rows
    case_rows += c_rows

deadlines_df = pd.DataFrame(deadline_rows)
cases_df = pd.DataFrame(case_rows)

combined = []
for _, cases in all_cases:
    combined.extend(cases)
client_label = all_cases[0][0] if len(all_cases) == 1 else "Combined"

# --- Summary metrics ------------------------------------------------------ #
m1, m2, m3 = st.columns(3)
m1.metric("Files", len(all_cases))
m2.metric("Cases", len(cases_df))
m3.metric("Deadlines in range", len(deadlines_df))

# --- Tabs ----------------------------------------------------------------- #
tab1, tab2, tab3 = st.tabs([f"\U0001F4CB Deadlines ({len(deadlines_df)})",
                            f"\U0001F5C2\uFE0F All Cases ({len(cases_df)})",
                            "\U0001F4C5 Calendar"])
with tab1:
    st.dataframe(deadlines_df, use_container_width=True, hide_index=True)
with tab2:
    st.dataframe(cases_df, use_container_width=True, hide_index=True)
with tab3:
    months = deadline_months(deadline_rows)
    if not months:
        st.info("No deadlines in range to display. Adjust the lookback slider above.")
    else:
        left, right = st.columns([3, 1])
        with left:
            choice = st.selectbox("Month", options=months, format_func=month_label)
        with right:
            st.write("")
            st.write("")
            pdf_bytes = month_pdf(deadline_rows, choice[0], choice[1],
                                  client_label=client_label)
            st.download_button(
                "\U0001F4C4 Download PDF",
                pdf_bytes,
                file_name=f"{client_label}_{choice[0]}-{choice[1]:02d}_calendar.pdf",
                mime="application/pdf",
                use_container_width=True,
            )
        st.markdown(month_grid_html(deadline_rows, choice[0], choice[1]),
                    unsafe_allow_html=True)

# --- Excel download ------------------------------------------------------- #
st.divider()
buf = BytesIO()
build_workbook(combined, client_label, deadline_cutoff=cutoff).save(buf)
buf.seek(0)
st.download_button(
    "\U0001F4E5 Download Excel workbook",
    buf,
    file_name=f"{client_label}_Docket_Extract.xlsx",
    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
)
